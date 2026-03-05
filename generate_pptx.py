#!/usr/bin/env python3
"""Generate Amogha Cafe & Restaurant product presentation (.pptx)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──
BURGUNDY = RGBColor(0x8B, 0x1A, 0x1A)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
GOLD_LIGHT = RGBColor(0xE8, 0xC5, 0x47)
DARK_BG = RGBColor(0x1E, 0x14, 0x0E)
DEEP_DARK = RGBColor(0x11, 0x0B, 0x06)
CREAM = RGBColor(0xFA, 0xF7, 0xF2)
WARM_GRAY = RGBColor(0x2C, 0x18, 0x10)
MUTED_TEXT = RGBColor(0xA0, 0x95, 0x88)
CARD_BG = RGBColor(0x2A, 0x1E, 0x16)
CARD_BORDER = RGBColor(0x3D, 0x2E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, left, top, width, height, font_size=18, color=CREAM,
             bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri', line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rich_text(slide, parts, left, top, width, height, font_size=18,
                  alignment=PP_ALIGN.LEFT, font_name='Calibri', line_spacing=1.2):
    """parts = [(text, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    for i, (text, color, bold) in enumerate(parts):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return txBox


def add_gold_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_card(slide, left, top, width, height, icon, title, body):
    card = add_shape_bg(slide, left, top, width, height, CARD_BG, CARD_BORDER)
    # Gold top accent line
    add_shape_bg(slide, left + Inches(0.05), top, width - Inches(0.1), Pt(2), GOLD)
    # Icon
    add_text(slide, icon, left + Inches(0.3), top + Inches(0.25), Inches(0.6), Inches(0.5),
             font_size=24, alignment=PP_ALIGN.LEFT)
    # Title
    add_text(slide, title, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), Inches(0.4),
             font_size=16, bold=True, color=CREAM, font_name='Calibri')
    # Body
    add_text(slide, body, left + Inches(0.3), top + Inches(1.05), width - Inches(0.6), height - Inches(1.3),
             font_size=11, color=MUTED_TEXT, line_spacing=1.5)
    return card


def add_stat(slide, left, top, number, label):
    add_text(slide, str(number), left, top, Inches(1.8), Inches(0.7),
             font_size=40, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_text(slide, label.upper(), left, top + Inches(0.65), Inches(1.8), Inches(0.3),
             font_size=8, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)


def add_bullet_list(slide, items, left, top, width, height, font_size=12, color=MUTED_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * 1.6)
    return txBox


def add_label_badge(slide, text, left, top):
    w, h = Inches(2.2), Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.background()
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    return shape


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / HERO
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DEEP_DARK)

# Subtle burgundy gradient circle (decorative)
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), Inches(-1), Inches(7), Inches(7))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x15, 0x0E, 0x09)
circle.line.fill.background()
circle.shadow.inherit = False

# Label
add_label_badge(slide, "PRODUCT  OVERVIEW  2026", Inches(5.55), Inches(1.3))

# Title
add_rich_text(slide, [
    ("Amogha", GOLD, True),
    (" Cafe\n& Restaurant", CREAM, True),
], Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.2),
    font_size=52, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Gold divider
add_gold_line(slide, Inches(5.9), Inches(4.15), Inches(1.5))

# Subtitle
add_text(slide,
    "A full-stack, AI-powered restaurant platform — from online ordering\nto kitchen ops, delivery, and analytics.",
    Inches(2.5), Inches(4.4), Inches(8.3), Inches(0.8),
    font_size=16, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

# Stats row
stats = [("9", "App Surfaces"), ("11", "AI Features"), ("16+", "API Endpoints"), ("2,059", "Tests")]
start_x = Inches(2.0)
for i, (num, lbl) in enumerate(stats):
    add_stat(slide, start_x + Inches(i * 2.5), Inches(5.5), num, lbl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE CHALLENGE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_label_badge(slide, "THE  CHALLENGE", Inches(5.55), Inches(0.6))

add_rich_text(slide, [
    ("Restaurants Run on ", CREAM, True),
    ("Disconnected Tools", GOLD, True),
], Inches(1.5), Inches(1.15), Inches(10.3), Inches(0.8),
    font_size=36, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_text(slide,
    "Separate apps for ordering, kitchen management, delivery, payments, and analytics\ncreate chaos, data silos, and poor customer experiences.",
    Inches(2.5), Inches(2.0), Inches(8.3), Inches(0.7),
    font_size=14, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

cards = [
    ("\U0001F4F1 \u2194 \U0001F4BB", "Fragmented Systems", "Staff juggles 5+ tools for daily operations — POS, kitchen tickets, delivery apps, accounting software."),
    ("\U0001F50D", "No Intelligence", "Menu decisions, demand forecasting, and customer insights require manual analysis with spreadsheets."),
    ("\U0001F6D2", "Lost Revenue", "Poor ordering UX, no personalization, no upsell — customers leave without completing orders."),
]
for i, (icon, title, body) in enumerate(cards):
    add_card(slide, Inches(0.8 + i * 4.15), Inches(3.2), Inches(3.8), Inches(3.5), icon, title, body)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ONE PLATFORM, EVERY SURFACE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_DARK)

add_label_badge(slide, "OUR  SOLUTION", Inches(5.55), Inches(0.5))

add_rich_text(slide, [
    ("One Platform. ", CREAM, True),
    ("Every Surface.", GOLD, True),
], Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.7),
    font_size=36, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_text(slide,
    "Amogha unifies the entire restaurant operation into a single real-time platform powered by AI.",
    Inches(2.5), Inches(1.7), Inches(8.3), Inches(0.5),
    font_size=14, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

surfaces = [
    ("\U0001F4F1", "Online Ordering", "PWA with smart search, AI chatbot, recommendations, Razorpay checkout."),
    ("\U0001F4FA", "Kitchen Display", "Real-time order queue with prep timers, priority sorting, station filtering."),
    ("\U0001F4B3", "POS Terminal", "Staff-facing counter ordering with receipt generation and quick-access menu."),
    ("\U0001F6F5", "Delivery App", "Driver assignments, navigation, earnings tracking, live status updates."),
    ("\U0001F916", "Self-Service Kiosk", "Touch-first ordering with voice support and multi-tenant capability."),
    ("\U0001F4CA", "Admin Dashboard", "Orders, menu, inventory, staff, expenses, and AI-powered analytics."),
]
for i, (icon, title, body) in enumerate(surfaces):
    col = i % 3
    row = i // 3
    add_card(slide, Inches(0.6 + col * 4.15), Inches(2.5 + row * 2.65), Inches(3.8), Inches(2.35),
             icon, title, body)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — AI FEATURES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_label_badge(slide, "POWERED  BY  GEMINI  2.0  FLASH", Inches(5.1), Inches(0.5))

add_rich_text(slide, [
    ("11", GOLD, True),
    (" AI-Powered Features", CREAM, True),
], Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.7),
    font_size=36, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_text(slide,
    "Every AI feature runs server-side via Google Vertex AI, protecting API keys and enabling rich context.",
    Inches(2.5), Inches(1.7), Inches(8.3), Inches(0.5),
    font_size=14, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

# Left card — Customer
left_card = add_shape_bg(slide, Inches(0.6), Inches(2.5), Inches(5.9), Inches(4.5), CARD_BG, CARD_BORDER)
add_shape_bg(slide, Inches(0.65), Inches(2.5), Inches(5.8), Pt(2), GOLD)
add_text(slide, "\U0001F4AC  Customer-Facing", Inches(1.0), Inches(2.75), Inches(5.0), Inches(0.4),
         font_size=18, bold=True, color=CREAM)
add_bullet_list(slide, [
    "Conversational AI chatbot for ordering & FAQs",
    "Smart semantic menu search (\"something spicy under \u20b9300\")",
    "Personalized meal & combo recommendations",
    "AI-generated subscription meal plans",
    "\"AI For You\" personalized section on homepage",
], Inches(1.0), Inches(3.3), Inches(5.2), Inches(3.5), font_size=13)

# Right card — Admin
right_card = add_shape_bg(slide, Inches(6.8), Inches(2.5), Inches(5.9), Inches(4.5), CARD_BG, CARD_BORDER)
add_shape_bg(slide, Inches(6.85), Inches(2.5), Inches(5.8), Pt(2), GOLD)
add_text(slide, "\U0001F4BC  Operations & Admin", Inches(7.2), Inches(2.75), Inches(5.0), Inches(0.4),
         font_size=18, bold=True, color=CREAM)
add_bullet_list(slide, [
    "Bill/receipt OCR with expense categorization",
    "Natural language analytics queries",
    "Demand forecasting & inventory planning",
    "Automated review sentiment analysis",
    "Menu profitability insights",
    "Smart push notification generation",
], Inches(7.2), Inches(3.3), Inches(5.2), Inches(3.5), font_size=13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CUSTOMER EXPERIENCE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_DARK)

add_label_badge(slide, "CUSTOMER  EXPERIENCE", Inches(5.35), Inches(0.5))

add_rich_text(slide, [
    ("Designed to ", CREAM, True),
    ("Delight & Convert", GOLD, True),
], Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.7),
    font_size=36, alignment=PP_ALIGN.CENTER, font_name='Calibri')

features = [
    ("\U0001F4B0", "Smart Payments", "UPI, cards, wallets via Razorpay. Coupons, gift cards, split bill, COD."),
    ("\u2B50", "Loyalty & Badges", "3-tier rewards (Bronze / Silver / Gold), 10 badges, birthday bonuses."),
    ("\U0001F46B", "Group Ordering", "Shared carts for friends & family with real-time sync and split payments."),
    ("\U0001F30F", "Multi-Language", "Full support for English, Hindi, and Telugu across every screen."),
    ("\U0001F4CD", "Live Tracking", "Real-time delivery map with Leaflet.js and OpenStreetMap integration."),
    ("\U0001F4C5", "Reservations", "Table booking with time slots and WhatsApp confirmation."),
    ("\U0001F514", "Push Notifications", "FCM-powered alerts for order status, promotions, and milestones."),
    ("\U0001F4E6", "Subscriptions", "AI-generated meal plans with recurring orders and preference learning."),
]
for i, (icon, title, body) in enumerate(features):
    col = i % 4
    row = i // 4
    add_card(slide, Inches(0.4 + col * 3.2), Inches(2.0 + row * 2.75), Inches(2.95), Inches(2.45),
             icon, title, body)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_label_badge(slide, "TECHNICAL  ARCHITECTURE", Inches(5.25), Inches(0.5))

add_rich_text(slide, [
    ("Built on ", CREAM, True),
    ("Firebase", GOLD, True),
], Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.7),
    font_size=36, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_text(slide,
    "Real-time data sync across all 9 surfaces with zero server management.",
    Inches(2.5), Inches(1.7), Inches(8.3), Inches(0.5),
    font_size=14, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

# Architecture diagram box
arch_bg = add_shape_bg(slide, Inches(2.5), Inches(2.4), Inches(8.3), Inches(3.0),
                       RGBColor(0x15, 0x0E, 0x09), CARD_BORDER)

arch_layers = [
    "9 App Surfaces  (PWA + Kiosk + POS + KDS + ...)",
    "\u2193",
    "Vite 5  \u2192  Bundled ES Modules  \u2192  Firebase Hosting",
    "\u2193",
    "Cloud Firestore  (17+ collections, real-time listeners)",
    "\u2193",
    "Cloud Functions  (Express API, 16+ endpoints)",
    "\u2193",
    "Vertex AI (Gemini 2.0)  +  Razorpay  +  FCM",
]
arch_text = "\n".join(arch_layers)
txBox = slide.shapes.add_textbox(Inches(3.0), Inches(2.5), Inches(7.3), Inches(2.8))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(arch_layers):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    if line == "\u2193":
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED_TEXT
        p.font.name = 'Consolas'
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    else:
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = GOLD
        p.font.bold = True
        p.font.name = 'Consolas'
        p.space_after = Pt(0)
        p.space_before = Pt(2)

# Tech badges row
badges = ["Vanilla JS", "Vite 5", "Firestore", "Cloud Functions", "Gemini 2.0",
          "Razorpay", "Leaflet.js", "Chart.js", "FCM", "Playwright", "Vitest/Jest", "GitHub Actions"]
badge_w = Inches(1.5)
badge_h = Inches(0.32)
cols = 6
for i, badge_text in enumerate(badges):
    col = i % cols
    row = i // cols
    x = Inches(0.9 + col * 2.05)
    y = Inches(5.7 + row * 0.45)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.85), badge_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = badge_text
    p.font.size = Pt(9)
    p.font.color.rgb = CREAM
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DATA MODEL
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_DARK)

add_label_badge(slide, "DATA  MODEL", Inches(5.55), Inches(0.5))

add_rich_text(slide, [
    ("17+", GOLD, True),
    (" Firestore Collections", CREAM, True),
], Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.7),
    font_size=36, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_gold_line(slide, Inches(5.9), Inches(1.75), Inches(1.5))

collections = [
    ("\U0001F6D2", "Commerce", ["menu", "orders", "addons", "coupons / giftCards", "specials"]),
    ("\U0001F464", "Users & Engagement", ["users", "reviews", "referrals", "chatHistory", "reservations"]),
    ("\u2699\uFE0F", "Operations", ["inventory", "staff / expenses", "deliveryPersons", "tables", "settings / branches"]),
]
for i, (icon, title, items) in enumerate(collections):
    x = Inches(0.6 + i * 4.15)
    card = add_shape_bg(slide, x, Inches(2.2), Inches(3.8), Inches(4.5), CARD_BG, CARD_BORDER)
    add_shape_bg(slide, x + Inches(0.05), Inches(2.2), Inches(3.7), Pt(2), GOLD)
    add_text(slide, f"{icon}  {title}", x + Inches(0.3), Inches(2.45), Inches(3.2), Inches(0.4),
             font_size=18, bold=True, color=CREAM)
    add_bullet_list(slide, items, x + Inches(0.3), Inches(3.1), Inches(3.2), Inches(3.2), font_size=14)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — QUALITY & SECURITY
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_label_badge(slide, "ENGINEERING  QUALITY", Inches(5.35), Inches(0.5))

add_rich_text(slide, [
    ("Tested. Secured. ", CREAM, True),
    ("Production-Ready.", GOLD, True),
], Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.7),
    font_size=36, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_gold_line(slide, Inches(5.9), Inches(1.75), Inches(1.5))

# Stats
q_stats = [("2,059", "Unit Tests"), ("40", "API Tests"), ("E2E", "Playwright"), ("CI/CD", "Auto-Deploy")]
for i, (num, lbl) in enumerate(q_stats):
    add_stat(slide, Inches(1.8 + i * 2.7), Inches(2.1), num, lbl)

# Security card
sec_card = add_shape_bg(slide, Inches(0.6), Inches(3.5), Inches(5.9), Inches(3.5), CARD_BG, CARD_BORDER)
add_shape_bg(slide, Inches(0.65), Inches(3.5), Inches(5.8), Pt(2), GOLD)
add_text(slide, "\U0001F6E1\uFE0F  Security", Inches(1.0), Inches(3.75), Inches(5.0), Inches(0.4),
         font_size=18, bold=True, color=CREAM)
add_bullet_list(slide, [
    "Admin API key authentication middleware",
    "Rate limiting (30 req/min) on AI endpoints",
    "Prompt injection sanitization (control chars, 2K limit)",
    "Firestore security rules for access control",
    "Input validation on all endpoints",
], Inches(1.0), Inches(4.3), Inches(5.2), Inches(2.5), font_size=13)

# Performance card
perf_card = add_shape_bg(slide, Inches(6.8), Inches(3.5), Inches(5.9), Inches(3.5), CARD_BG, CARD_BORDER)
add_shape_bg(slide, Inches(6.85), Inches(3.5), Inches(5.8), Pt(2), GOLD)
add_text(slide, "\u26A1  Performance", Inches(7.2), Inches(3.75), Inches(5.0), Inches(0.4),
         font_size=18, bold=True, color=CREAM)
add_bullet_list(slide, [
    "Lazy module loading via requestIdleCallback",
    "In-memory menu cache (10-min TTL) in Cloud Functions",
    "Service workers for offline capability",
    "Vite-bundled ES modules with tree-shaking",
    "Deferred AI recommendations (4s delay)",
], Inches(7.2), Inches(4.3), Inches(5.2), Inches(2.5), font_size=13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LIVE IN PRODUCTION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_DARK)

add_label_badge(slide, "LIVE  IN  PRODUCTION", Inches(5.35), Inches(0.5))

add_rich_text(slide, [
    ("Amogha", GOLD, True),
    (" — Hyderabad", CREAM, True),
], Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.7),
    font_size=36, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_gold_line(slide, Inches(5.9), Inches(1.75), Inches(1.5))

# Restaurant card
rest_card = add_shape_bg(slide, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5), CARD_BG, CARD_BORDER)
add_shape_bg(slide, Inches(0.85), Inches(2.2), Inches(5.5), Pt(2), GOLD)
add_text(slide, "\U0001F35B", Inches(2.8), Inches(2.5), Inches(1.5), Inches(0.6),
         font_size=36, alignment=PP_ALIGN.CENTER)
add_text(slide, "Authentic Indian Cuisine", Inches(1.2), Inches(3.1), Inches(5.0), Inches(0.4),
         font_size=20, bold=True, color=CREAM, alignment=PP_ALIGN.CENTER)
add_text(slide,
    "South Indian, Hyderabadi, and Pan-Indian\nspecialties served from Kukatpally, Hyderabad.",
    Inches(1.2), Inches(3.6), Inches(5.0), Inches(0.7),
    font_size=13, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.5)
add_text(slide,
    "\U0001F4CD Pragathi Nagar Rd, Kukatpally\n\U0001F4DE +91 91210 04999\n\u23F0 Mon-Thu 11:00-21:30 \u00b7 Fri-Sat till 22:30",
    Inches(1.2), Inches(4.6), Inches(5.0), Inches(1.2),
    font_size=12, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.6)

# Reviews card
rev_card = add_shape_bg(slide, Inches(6.9), Inches(2.2), Inches(5.6), Inches(4.5), CARD_BG, CARD_BORDER)
add_shape_bg(slide, Inches(6.95), Inches(2.2), Inches(5.5), Pt(2), GOLD)
add_text(slide, "\u2B50", Inches(8.9), Inches(2.5), Inches(1.5), Inches(0.6),
         font_size=36, alignment=PP_ALIGN.CENTER)
add_text(slide, "Customer Love", Inches(7.3), Inches(3.1), Inches(5.0), Inches(0.4),
         font_size=20, bold=True, color=CREAM, alignment=PP_ALIGN.CENTER)

add_stat(slide, Inches(7.5), Inches(3.7), "4.8", "Rating")
add_stat(slide, Inches(10.0), Inches(3.7), "1,250+", "Reviews")

add_text(slide,
    "Free delivery on orders above \u20b9500.\nHappy Hour deals daily.\n3-tier loyalty program.",
    Inches(7.3), Inches(5.1), Inches(5.0), Inches(1.0),
    font_size=13, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_DARK)

# Decorative circle
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), Inches(0), Inches(7), Inches(7))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x15, 0x0E, 0x09)
circle.line.fill.background()
circle.shadow.inherit = False

add_rich_text(slide, [
    ("The ", CREAM, True),
    ("Complete", GOLD, True),
    (" Restaurant\nPlatform.", CREAM, True),
], Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.8),
    font_size=48, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_gold_line(slide, Inches(5.9), Inches(3.7), Inches(1.5))

add_text(slide,
    "One codebase. Nine surfaces. Real-time everything.\nPowered by AI. Built on Firebase. Ready to scale.",
    Inches(2.5), Inches(4.0), Inches(8.3), Inches(0.8),
    font_size=16, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.6)

# URL badges
for i, url in enumerate(["amoghahotels.com", "amogha-cafe.web.app"]):
    x = Inches(4.9 + i * 2.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.1), Inches(2.1), Inches(0.4))
    shape.fill.background()
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = url
    p.font.size = Pt(11)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

add_text(slide,
    "Built with Firebase \u00b7 Vertex AI \u00b7 Razorpay \u00b7 Vite\n\u00a9 2026 Amogha Cafe & Restaurant",
    Inches(3.5), Inches(6.0), Inches(6.3), Inches(0.7),
    font_size=10, color=RGBColor(0x60, 0x55, 0x4A), alignment=PP_ALIGN.CENTER, line_spacing=1.6)


# ── Save ──
output_path = '/home/user/amogha-cafe/Amogha_Cafe_Presentation.pptx'
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
